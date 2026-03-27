"""RAG 문서 적재 품질을 높이기 위한 공통 유틸리티입니다.

PDF, DOCX, PPTX 텍스트를 정리하고,
문단과 조항 경계를 최대한 살린 청크로 변환해 벡터DB에 넣습니다.
"""

from __future__ import annotations

import os
import re
from pathlib import Path

import chromadb
import docx
import pdfplumber
from pptx import Presentation
from sentence_transformers import SentenceTransformer

from core.settings import get_rag_collection_name, get_rag_embedding_model


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}
HEADING_PATTERN = re.compile(r"^(제\s*\d+\s*조|\d+\s*[\.\)]|[0-9]+\s+[가-힣A-Za-z].*)")
MAX_CHUNK_LENGTH = 900
MIN_CHUNK_LENGTH = 180
CHUNK_OVERLAP = 120
MODEL_NAME = get_rag_embedding_model()


def normalize_text(text: str) -> str:
    """추출된 원문을 검색에 유리한 형태로 정리합니다."""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\u00a0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"(?<!\n)-\s+", "-", text)
    return text.strip()


def read_pdf(path: Path) -> str:
    """PDF 파일의 텍스트를 페이지 순서대로 읽습니다."""
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                pages.append(page_text)
    return normalize_text("\n\n".join(pages))


def read_docx(path: Path) -> str:
    """DOCX 문서의 문단 텍스트를 읽습니다."""
    document = docx.Document(path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    return normalize_text("\n\n".join(paragraphs))


def read_pptx(path: Path) -> str:
    """PPTX 슬라이드의 텍스트를 읽습니다."""
    slide_texts = []
    presentation = Presentation(path)
    for slide in presentation.slides:
        shape_texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                shape_texts.append(text.strip())
        if shape_texts:
            slide_texts.append("\n".join(shape_texts))
    return normalize_text("\n\n".join(slide_texts))


def read_document_text(path: Path) -> str:
    """확장자에 맞는 리더로 문서 본문을 추출합니다."""
    extension = path.suffix.lower()
    if extension == ".pdf":
        return read_pdf(path)
    if extension == ".docx":
        return read_docx(path)
    if extension == ".pptx":
        return read_pptx(path)
    raise ValueError(f"지원하지 않는 문서 형식입니다: {path.name}")


def split_into_sections(text: str):
    """문단과 조항 제목을 기준으로 섹션을 나눕니다."""
    paragraphs = [paragraph.strip() for paragraph in text.split("\n\n") if paragraph.strip()]
    if not paragraphs:
        return []

    sections = []
    current_title = "문서 본문"
    current_parts = []

    for paragraph in paragraphs:
        first_line = paragraph.splitlines()[0].strip()
        if HEADING_PATTERN.match(first_line) and current_parts:
            sections.append((current_title, "\n\n".join(current_parts)))
            current_title = first_line
            current_parts = [paragraph]
            continue

        if not current_parts:
            current_title = first_line if len(first_line) <= 80 else current_title
        current_parts.append(paragraph)

    if current_parts:
        sections.append((current_title, "\n\n".join(current_parts)))
    return sections


def _split_long_text_with_overlap(text: str, max_length: int = MAX_CHUNK_LENGTH, overlap: int = CHUNK_OVERLAP):
    """긴 섹션을 문자 길이 기준으로 겹치게 나눕니다."""
    chunks = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(text_length, start + max_length)
        if end < text_length:
            split_index = text.rfind("\n\n", start, end)
            if split_index == -1:
                split_index = text.rfind(". ", start, end)
            if split_index != -1 and split_index > start + MIN_CHUNK_LENGTH:
                end = split_index + (2 if text[split_index:split_index + 2] == ". " else 0)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= text_length:
            break
        start = max(end - overlap, start + 1)

    return chunks


def build_document_chunks(filename: str, text: str):
    """문서 내용을 제목/문단 경계를 살려 청크 목록으로 변환합니다."""
    sections = split_into_sections(text)
    chunks = []

    for section_title, section_text in sections:
        section_text = normalize_text(section_text)
        if not section_text:
            continue

        section_chunks = _split_long_text_with_overlap(section_text)
        for section_chunk in section_chunks:
            cleaned_chunk = section_chunk.strip()
            if len(cleaned_chunk) < 40:
                continue
            chunk_index = len(chunks)
            chunk_text = f"[문서명] {Path(filename).stem}\n[섹션] {section_title}\n\n{cleaned_chunk}"
            chunks.append(
                {
                    "id": f"{filename}_{chunk_index:04d}",
                    "text": chunk_text,
                    "metadata": {
                        "filename": filename,
                        "section_title": section_title,
                        "chunk_index": chunk_index,
                    },
                }
            )

    return chunks


def iter_document_paths(doc_dir: Path):
    """지원하는 문서 파일 경로를 이름순으로 반환합니다."""
    for path in sorted(doc_dir.iterdir()):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path


def build_file_signature(path: Path):
    """파일 변경 감지를 위한 간단한 시그니처를 생성합니다."""
    stat = path.stat()
    return (int(stat.st_mtime), stat.st_size)


def create_embedding_model(cache_dir: str):
    """임베딩 모델을 공통 옵션으로 생성합니다."""
    return SentenceTransformer(MODEL_NAME, cache_folder=cache_dir)


def create_collection(vector_db_dir: str):
    """Jarvis 문서 컬렉션을 열거나 생성합니다."""
    client = chromadb.PersistentClient(path=vector_db_dir)
    return client.get_or_create_collection(name=get_rag_collection_name(), metadata={"hnsw:space": "cosine"})


def remove_document_from_collection(collection, filename: str):
    """지정한 파일명에 해당하는 기존 청크를 삭제합니다."""
    existing = collection.get(where={"filename": filename}, include=["metadatas"])
    ids = existing.get("ids", [])
    if ids:
        collection.delete(ids=ids)


def upsert_document(collection, embedding_model, filename: str, text: str):
    """문서 한 개를 다시 청크화하고 벡터DB에 반영합니다."""
    chunks = build_document_chunks(filename, text)
    if not chunks:
        remove_document_from_collection(collection, filename)
        return 0

    remove_document_from_collection(collection, filename)
    texts = [chunk["text"] for chunk in chunks]
    embeddings = embedding_model.encode(texts).tolist()
    collection.add(
        ids=[chunk["id"] for chunk in chunks],
        embeddings=embeddings,
        documents=texts,
        metadatas=[chunk["metadata"] for chunk in chunks],
    )
    return len(chunks)